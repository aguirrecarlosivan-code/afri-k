import os
import logging
from typing import Dict, Any, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger("afrik.reports.pptx")


class PPTXReportGenerator:
    """
    Generates professional 4-slide PowerPoint Executive Pitch Deck for Afri-k (Once Noticias).
    """

    @staticmethod
    def generate_presentation(
        output_path: str,
        report_data: Dict[str, Any],
        ai_analysis: Dict[str, Any],
        platform_summaries: List[Dict[str, Any]],
    ) -> str:
        """
        Creates PowerPoint .pptx presentation at output_path.
        """
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 Widescreen
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]

        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout)
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        title.text = "Afri-k: Informe Ejecutivo de Once Noticias"
        subtitle.text = f"Análisis de Inteligencia Editorial | Período: {report_data.get('period_start', '')} - {report_data.get('period_end', '')}"

        # Slide 2: Resumen Ejecutivo & KPIs
        slide2 = prs.slides.add_slide(blank_layout)
        txBox = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "1. Resumen Ejecutivo & Alcance Global"
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(15, 23, 42)

        txBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = ai_analysis.get("executive_summary", "")
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(51, 65, 85)

        # Slide 3: Recomendaciones & Sentimiento
        slide3 = prs.slides.add_slide(blank_layout)
        txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = "2. Recomendaciones Estratégicas & Sentimiento de Audiencia"
        p3.font.bold = True
        p3.font.size = Pt(24)

        txBox4 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
        tf4 = txBox4.text_frame
        for idx, rec in enumerate(ai_analysis.get("recommendations", []), start=1):
            p_rec = tf4.add_paragraph() if idx > 1 else tf4.paragraphs[0]
            p_rec.text = f"{idx}. {rec}"
            p_rec.font.size = Pt(16)
            p_rec.space_after = Pt(12)

        # Slide 4: Métricas por Plataforma
        slide4 = prs.slides.add_slide(blank_layout)
        txBox5 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf5 = txBox5.text_frame
        p5 = tf5.paragraphs[0]
        p5.text = "3. Desglose de Desempeño por Canal"
        p5.font.bold = True
        p5.font.size = Pt(24)

        txBox6 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
        tf6 = tf5.text_frame if hasattr(tf5, "text_frame") else txBox6.text_frame
        for ps in platform_summaries:
            p_ps = tf6.add_paragraph()
            p_ps.text = f"• {ps.get('platform', '').capitalize()}: {ps.get('followers', 0):,} seguidores | Alcance: {ps.get('total_reach', 0):,} | Engagement: {ps.get('avg_engagement', 0.0)}%"
            p_ps.font.size = Pt(15)
            p_ps.space_after = Pt(8)

        prs.save(output_path)
        logger.info(f"PowerPoint Executive Presentation successfully generated at {output_path}")
        return output_path
